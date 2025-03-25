"""
PowerPoint Export Callbacks Module

This module contains callbacks for generating and downloading PowerPoint presentations
with visualizations and data for selected countries.
"""

import os
import io
import pandas as pd
import plotly.io as pio
from datetime import datetime
from dash import Input, Output, State, dcc
from dash.exceptions import PreventUpdate

# PowerPoint libraries
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Import configuration
from ..config_and_data import (
    config, 
    df_variables, 
    df_names, 
    df_index,
    get_dashboard_file_path
)

from ..utils import order_dataframe_by_country


def register_ppt_callbacks(app, cache):
    """
    Register PowerPoint export-related callbacks with the app.
    
    Args:
        app: The Dash application instance
        cache: The Flask-Cache instance
    """
    
    @app.callback(
        Output("download-ppt", "data"),
        [Input("download-ppt-button", "n_clicks")],
        [State("extract-country-checklist", "value"),
         State("selected-metric", "data"),
         State("selected-scenario", "data")]
    )
    def generate_ppt(n_clicks, selected_countries, metric, scenario):
        """
        Generate a PowerPoint presentation with visualizations for the selected countries.
        
        Args:
            n_clicks: Number of times the download button was clicked
            selected_countries: List of selected countries
            metric: Selected metric filter
            scenario: Selected scenario
            
        Returns:
            dict: Download data for the PowerPoint file
        """
        # Don't generate if download button not clicked
        if not n_clicks:
            raise PreventUpdate

        # If "All" is selected, use all available countries
        if selected_countries and "All" in selected_countries:
            selected_countries = sorted(df_names["Country"].dropna().unique().tolist())
        
        # Don't generate if no countries selected
        if not selected_countries:
            raise PreventUpdate

        # Import necessary callbacks for plot generation
        from .callback_plots import (
            create_polar_plot, 
            create_bar_plot, 
            create_financial_chart, 
            create_multiple_box_plot)
        # Get dashboard path
        dashboard_file = get_dashboard_file_path()

        # Get PowerPoint configuration
        ppt_cfg = config["ppt"]["config"]
        
        # Create a presentation using the theme file specified in config
        ppt_theme = config["file_paths"]["ppt_theme"]
        prs = Presentation(ppt_theme)

        # --- COVER SLIDE ---
        cover_layout = prs.slide_layouts[ppt_cfg["cover_layout_index"]]
        cover_slide = prs.slides.add_slide(cover_layout)
        
        # Set the cover title
        if cover_slide.shapes.title:
            cover_slide.shapes.title.text = ppt_cfg["cover_title"]
        
        # Format date for cover slide
        formatted_date = datetime.now().strftime(ppt_cfg["ppt_date_format"])
        
        # Set the cover date
        if len(cover_slide.placeholders) > 1:
            cover_slide.placeholders[1].text = f"Date: {formatted_date}"

        # --- TABLE SLIDES: Split the FRV table into 3 parts ---
        # Load dashboard data
        df_dash = pd.read_excel(dashboard_file, sheet_name="Dashboard")
        
        if len(selected_countries) > 0:
            ordered_df = order_dataframe_by_country(df_dash, selected_countries[0])
        else:
            ordered_df = df_dash

        # Get column configurations for each metric type
        tables_cfg = ppt_cfg["tables_config"]
        fundamentals_cols = tables_cfg["fundamentals"]
        valuations_cols = tables_cfg["valuations"]
        risks_cols = tables_cfg["risks"]

        # Define valid column sets for each metric type
        from ..config_and_data import valid_fund, valid_risk, valid_val

        def build_metric_table(df, base_cols, valid_set):
            """Build a table with base columns and additional columns from a valid set."""
            cols = list(base_cols)
            if "TotalScore" not in cols:
                cols.append("TotalScore")
            additional = [col for col in df.columns if col in valid_set]
            for col in additional:
                if col not in cols:
                    cols.append(col)
            return df[cols].copy()

        # Build tables for each metric type
        df_fund = build_metric_table(ordered_df, fundamentals_cols, valid_fund)
        df_val = build_metric_table(ordered_df, valuations_cols, valid_val)
        df_risk = build_metric_table(ordered_df, risks_cols, valid_risk)

        # Sort tables by their respective scores
        if "Fundamental Score" in df_fund.columns:
            df_fund = df_fund.sort_values(by="Fundamental Score", ascending=False)
        if "Valuation Score" in df_val.columns:
            df_val = df_val.sort_values(by="Valuation Score", ascending=False)
        if "Risk Score" in df_risk.columns:
            df_risk = df_risk.sort_values(by="Risk Score", ascending=False)

        # Add slides for the three tables
        add_table_slide("Fundamentals Table", df_fund, prs, ppt_cfg)
        add_table_slide("Valuations Table", df_val, prs, ppt_cfg)
        add_table_slide("Risks Table", df_risk, prs, ppt_cfg)

        # --- OVERVIEW SLIDE: "Total Shareholder Return" / "Regional Overview" ---
        overview_layout = prs.slide_layouts[ppt_cfg["graph_layout_index"]]
        overview_slide = prs.slides.add_slide(overview_layout)
        
        # Set overview titles
        overview_title = ppt_cfg["overview"]["title"]
        overview_sub = ppt_cfg["overview"]["subtitle"]
        
        if overview_slide.shapes.title:
            overview_slide.shapes.title.text = overview_title
            
        if len(overview_slide.placeholders) > 1:
            overview_slide.placeholders[13].text = overview_sub
            
        # Import stacked bar chart function
        from ..layout import make_stacked_bar_for_all_indices
            
        # Generate and add stacked bar chart
        stacked_graph = make_stacked_bar_for_all_indices()
        stacked_img_bytes = pio.to_image(stacked_graph.figure, format="png")
        
        # Add the chart to the slide
        placeholders_ov = {ph.placeholder_format.idx: ph for ph in overview_slide.placeholders}
        if 15 in placeholders_ov:
            ph = placeholders_ov[15]
            overview_slide.shapes.add_picture(
                io.BytesIO(stacked_img_bytes),
                ph.left, ph.top, 
                width=ph.width, 
                height=ph.height
            )
            # Remove the used placeholder
            try:
                ph.element.getparent().remove(ph.element)
            except Exception as e:
                print(f"Error removing overview placeholder: {e}")

        # --- SLIDES FOR EACH COUNTRY ---
        for country in selected_countries:
            # --- GRAPH SLIDE: Polar and Bar charts ---
            graph_layout = prs.slide_layouts[ppt_cfg["graph_layout_index"]]
            graph_slide = prs.slides.add_slide(graph_layout)
            
            # Set slide title and subtitle
            if graph_slide.shapes.title:
                graph_slide.shapes.title.text = f"{country}"
                
            if len(graph_slide.placeholders) > 1:
                graph_slide.placeholders[13].text = ppt_cfg["graph_slide_subtitle"]
                
            # Generate polar and bar plots
            polar_fig = create_polar_plot(country, metric)
            bar_fig = create_bar_plot(country, scenario)
            
            # Convert plots to images
            polar_img_bytes = pio.to_image(polar_fig, format="png")
            bar_img_bytes = pio.to_image(bar_fig, format="png")
            
            # Add plots to the slide
            placeholders_graph = {ph.placeholder_format.idx: ph for ph in graph_slide.placeholders}
            
            # Add polar plot
            if 15 in placeholders_graph:
                ph = placeholders_graph[15]
                graph_slide.shapes.add_picture(
                    io.BytesIO(polar_img_bytes),
                    ph.left, ph.top, 
                    width=ph.width, 
                    height=ph.height
                )
                try:
                    ph.element.getparent().remove(ph.element)
                except Exception as e:
                    print(f"Error removing graph placeholder 15: {e}")
                    
            # Add bar plot
            if 16 in placeholders_graph:
                ph = placeholders_graph[16]
                graph_slide.shapes.add_picture(
                    io.BytesIO(bar_img_bytes),
                    ph.left, ph.top, 
                    width=ph.width, 
                    height=ph.height
                )
                try:
                    ph.element.getparent().remove(ph.element)
                except Exception as e:
                    print(f"Error removing graph placeholder 16: {e}")
        
            # --- FINANCIAL SLIDE: Index comparison and box plots ---
            fin_layout = prs.slide_layouts[ppt_cfg["financial_layout_index"]]
            fin_slide = prs.slides.add_slide(fin_layout)
            
            # Set slide title and subtitle
            if fin_slide.shapes.title:
                fin_slide.shapes.title.text = f"{country}"
                
            if len(fin_slide.placeholders) > 1:
                fin_slide.placeholders[13].text = ppt_cfg["financial_slide_subtitle"]
                
            # Get placeholders for financial slide
            placeholders_fin = {ph.placeholder_format.idx: ph for ph in fin_slide.placeholders}
            
            # Generate financial charts
            index_comp_fig = create_financial_chart(country)
            multiple_fig = create_multiple_box_plot(country)
            
            # Convert charts to images
            index_img_bytes = pio.to_image(index_comp_fig, format="png")
            multiple_img_bytes = pio.to_image(multiple_fig, format="png")
            
            # Add index comparison chart
            if 16 in placeholders_fin:
                ph = placeholders_fin[16]
                fin_slide.shapes.add_picture(
                    io.BytesIO(index_img_bytes),
                    ph.left, ph.top, 
                    width=ph.width, 
                    height=ph.height
                )
                try:
                    ph.element.getparent().remove(ph.element)
                except Exception as e:
                    print(f"Error removing financial placeholder 16: {e}")
                    
            # Add box plots
            if 17 in placeholders_fin:
                ph = placeholders_fin[17]
                fin_slide.shapes.add_picture(
                    io.BytesIO(multiple_img_bytes),
                    ph.left, ph.top, 
                    width=ph.width, 
                    height=ph.height
                )
                try:
                    ph.element.getparent().remove(ph.element)
                except Exception as e:
                    print(f"Error removing financial placeholder 17: {e}")
        
        # --- FOOTER AND SLIDE NUMBERING ---
        # Format the date for the footer
        footer_date = datetime.now().strftime(ppt_cfg["ppt_filename_date_format"])
        
        # Get footer configuration
        footer_cfg = ppt_cfg["footer"]
        footer_text = footer_cfg["text"]
        footer_font_size = footer_cfg["font_size"]
        footer_color = footer_cfg["color"]
        
        # Add footer and slide numbers to all non-cover slides
        non_cover_slide_number = 1
        for slide in prs.slides:
            # Skip cover slide
            if slide.slide_layout == prs.slide_layouts[ppt_cfg["cover_layout_index"]]:
                continue
                
            # Process placeholders
            for ph in slide.placeholders:
                # Add footer text
                if ph.placeholder_format.idx == ppt_cfg["footer_placeholder_idx"]:
                    ph.text = f"{footer_text} | {footer_date}"
                    if ph.text_frame:
                        for paragraph in ph.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(footer_font_size)
                                run.font.color.rgb = RGBColor(*footer_color)
                                
                # Add slide number
                elif ph.placeholder_format.idx == ppt_cfg["slide_number_placeholder_idx"]:
                    ph.text = f"Slide {non_cover_slide_number}"
                    
            non_cover_slide_number += 1

        # Save the presentation to a BytesIO object
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        # Return the file for download
        return dcc.send_bytes(
            ppt_stream.read(), 
            filename=f"MacroExtraction_{footer_date}.pptx"
        )


def add_table_slide(slide_title, df_table, prs, ppt_cfg):
    """
    Create a slide with a data table using the style defined in the configuration.
    
    Args:
        slide_title: Title for the slide
        df_table: DataFrame containing the table data
        prs: PowerPoint presentation object
        ppt_cfg: PowerPoint configuration dictionary
    """
    # Get styles for header and data from configuration
    table_header_style = ppt_cfg["table_styles"]["header"]
    table_data_style = ppt_cfg["table_styles"]["data"]

    # Create the slide
    table_layout = prs.slide_layouts[ppt_cfg["table_layout_index"]]
    slide = prs.slides.add_slide(table_layout)

    # Set the slide title
    if slide.shapes.title:
        slide.shapes.title.text = slide_title

    # Set the subtitle if placeholder exists
    if len(slide.placeholders) > 1:
        slide.placeholders[13].text = ppt_cfg["table_slide_title"]

    # Define table position and dimensions
    left = Cm(1)
    top = Cm(3)
    width = Cm(23.1)
    height = Cm(8.98)

    # Get table dimensions from DataFrame
    rows, cols = df_table.shape
    
    # Create the table
    table_obj = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    # Set row heights
    table_obj.rows[0].height = Cm(table_header_style["first_col_width_cm"])  # Header row height
    for i in range(rows):
        table_obj.rows[i + 1].height = Cm(table_data_style["first_col_width_cm"])  # Data rows height

    # Format header row
    for j, col in enumerate(df_table.columns):
        cell = table_obj.cell(0, j)
        cell.text = str(col)
        
        # Format text in the cell
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(table_header_style["first_col_font_size"])
                run.font.color.rgb = RGBColor(*table_header_style["font_color"])
                
            # Set text alignment
            align = table_header_style["alignment"].lower()
            if align == "center":
                paragraph.alignment = PP_ALIGN.CENTER
            elif align == "right":
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.LEFT
                
        # Set cell background color
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*table_header_style["bg_color"])

    # Format data rows
    for i in range(rows):
        for j in range(cols):
            cell = table_obj.cell(i + 1, j)
            
            # Format cell value
            val = df_table.iloc[i, j]
            try:
                numeric_value = int(float(val))
                cell_text = str(numeric_value)
            except (ValueError, TypeError):
                cell_text = str(val)
                
            cell.text = cell_text
            
            # Format text in the cell
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Use different font size for first column vs other columns
                    if j == 0:
                        run.font.size = Pt(table_data_style["first_col_font_size"])
                    else:
                        run.font.size = Pt(table_data_style["other_cols_font_size"])
                        
                    run.font.color.rgb = RGBColor(*table_data_style["font_color"])
                    
                # Set text alignment
                align = table_data_style["alignment"].lower()
                if align == "center":
                    paragraph.alignment = PP_ALIGN.CENTER
                elif align == "right":
                    paragraph.alignment = PP_ALIGN.RIGHT
                else:
                    paragraph.alignment = PP_ALIGN.LEFT
                    
            # Set cell background color
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*table_data_style["bg_color"])